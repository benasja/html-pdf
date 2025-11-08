import os
import sys
import http.server
import socketserver
import socket
import subprocess
import threading
import traceback
import tempfile
import webbrowser
import re
from typing import Optional, Tuple

import customtkinter as ctk
from tkinter import filedialog, messagebox


def _ensure_playwright_browsers() -> None:
    """Ensure Playwright Chromium is installed and discoverable at runtime.

    In PyInstaller bundles, Playwright may default to a package-local browsers
    directory that isn't shipped. We redirect to the user cache path and
    install Chromium there if not already present.
    """
    # Use the default user cache for Playwright browsers (cross-platform)
    if sys.platform == "win32":
        user_cache = os.path.join(os.environ.get("LOCALAPPDATA", ""), "ms-playwright")
    elif sys.platform == "darwin":
        user_cache = os.path.expanduser("~/Library/Caches/ms-playwright")
    else:
        user_cache = os.path.expanduser("~/.cache/ms-playwright")
    os.environ.setdefault("PLAYWRIGHT_BROWSERS_PATH", user_cache)

    chromium_installed_marker = os.path.join(user_cache, "chromium-*")
    # Best-effort detection; if directory not present, attempt install
    if not os.path.isdir(user_cache) or not any(
        name.startswith("chromium") for name in os.listdir(user_cache)
    ):
        try:
            subprocess.run(
                [sys.executable, "-m", "playwright", "install", "chromium"],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
            )
        except Exception:
            # Let Playwright raise a clearer error later if install truly failed
            pass


def convert_html_to_pdf_sync(html_content: str, output_pdf_path: str, continuous: bool = False) -> None:
    """Render HTML to PDF using Playwright (Chromium) synchronously.

    Args:
        html_content: The complete HTML string to render.
        output_pdf_path: Absolute path to write the resulting PDF file.
    """
    from playwright.sync_api import sync_playwright  # Imported here to start fast UI

    _ensure_playwright_browsers()
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context()
        page = context.new_page()

        # Use screen media; wait for network to be idle so external CSS/images load
        page.emulate_media(media="screen")
        page.set_content(html_content, wait_until="networkidle")

        if continuous:
            # Measure full content size and generate a single tall page
            # Use CSS pixels; Chromium treats 1px = 1/96 inch
            size = page.evaluate(
                """
(() => {
  const el = document.documentElement;
  const body = document.body;
  const width = Math.max(el.scrollWidth, el.offsetWidth, body?.scrollWidth||0, body?.offsetWidth||0);
  const height = Math.max(el.scrollHeight, el.offsetHeight, body?.scrollHeight||0, body?.offsetHeight||0);
  return { width, height };
})()
                """
            )
            width_px = max(1, int(size["width"]))
            height_px = max(1, int(size["height"]))

            pdf_bytes = page.pdf(
                width=f"{width_px}px",
                height=f"{height_px}px",
                print_background=True,
                margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
                prefer_css_page_size=False,
            )
        else:
            pdf_bytes = page.pdf(
                format="A4",
                print_background=True,
                prefer_css_page_size=True,
            )

        with open(output_pdf_path, "wb") as f:
            f.write(pdf_bytes)

        context.close()
        browser.close()


def convert_html_to_png_sync(html_content: str, output_png_path: str) -> None:
    """Render HTML to a full-page PNG using Playwright (Chromium)."""
    from playwright.sync_api import sync_playwright

    _ensure_playwright_browsers()
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context()
        page = context.new_page()

        page.emulate_media(media="screen")
        page.set_content(html_content, wait_until="networkidle")

        png_bytes = page.screenshot(full_page=True, type="png")

        with open(output_png_path, "wb") as f:
            f.write(png_bytes)

        context.close()
        browser.close()


def convert_html_to_docx_sync(html_content: str, output_docx_path: str) -> None:
    """Convert HTML to DOCX by rasterizing to PNG and embedding it."""
    import tempfile
    from docx import Document
    from docx.shared import Inches

    with tempfile.TemporaryDirectory() as tmpdir:
        png_path = os.path.join(tmpdir, "page.png")
        convert_html_to_png_sync(html_content, png_path)

        doc = Document()
        # Fit image to typical page width; python-docx will keep aspect ratio
        doc.add_picture(png_path, width=Inches(6.5))
        doc.save(output_docx_path)


def convert_html_to_pptx_sync(html_content: str, output_pptx_path: str) -> None:
    """Convert HTML to PPTX creating one slide per .slide section if present.

    Fallback: if no .slide sections found, capture full page as a single slide.
    """
    import tempfile
    from pptx import Presentation
    from playwright.sync_api import sync_playwright

    _ensure_playwright_browsers()
    with tempfile.TemporaryDirectory() as tmpdir:
        screenshots: list[str] = []
        first_slide_ratio: Optional[float] = None

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            # Use a 16:9 viewport; element screenshots ignore viewport size for clipping,
            # but 100vh/100vw-based layouts will be consistent
            context = browser.new_context(viewport={"width": 1920, "height": 1080})
            page = context.new_page()

            page.emulate_media(media="screen")
            page.set_content(html_content, wait_until="networkidle")

            # Prefer <section class="slide">, else any .slide
            locator = page.locator("section.slide, .slide")
            count = locator.count()

            if count == 0:
                # Fallback: single screenshot of the full page
                png_path = os.path.join(tmpdir, "slide-1.png")
                page.screenshot(path=png_path, full_page=True, type="png")
                screenshots.append(png_path)
            else:
                for i in range(count):
                    # Element screenshots auto-scroll into view
                    el = locator.nth(i)
                    # Record aspect ratio (h/w) from first slide for PPTX slide sizing
                    if i == 0:
                        box = el.bounding_box()
                        if box and box.get("width") and box.get("height"):
                            first_slide_ratio = max(0.01, float(box["height"]) / float(box["width"]))
                    png_path = os.path.join(tmpdir, f"slide-{i+1}.png")
                    el.screenshot(path=png_path, type="png")
                    screenshots.append(png_path)

            context.close()
            browser.close()

        prs = Presentation()
        # If we detected a slide aspect ratio, size the PPTX accordingly to avoid
        # top/bottom whitespace when fitting images.
        if first_slide_ratio is not None:
            from pptx.util import Inches
            base_width_in = 13.333  # typical widescreen width
            prs.slide_width = Inches(base_width_in)
            prs.slide_height = Inches(base_width_in * first_slide_ratio)
        blank_layout = prs.slide_layouts[6]
        from PIL import Image
        for png in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide_w = prs.slide_width
            slide_h = prs.slide_height

            with Image.open(png) as im:
                img_w_px, img_h_px = im.size

            # Fit image within slide (contain), preserving aspect ratio
            target_w = slide_w
            target_h = int(slide_w * img_h_px / img_w_px)
            if target_h > slide_h:
                target_h = slide_h
                target_w = int(slide_h * img_w_px / img_h_px)

            left = int((slide_w - target_w) / 2)
            top = int((slide_h - target_h) / 2)

            slide.shapes.add_picture(png, left=left, top=top, width=target_w, height=target_h)
        prs.save(output_pptx_path)


class HtmlToPdfApp(ctk.CTk):
    def __init__(self) -> None:
        super().__init__()

        self.title("HTML-to-PDF Converter")
        self.geometry("900x700")
        self.minsize(700, 500)

        ctk.set_appearance_mode("System")
        ctk.set_default_color_theme("blue")

        # Layout configuration
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=0)

        # Text input for HTML
        self.html_text = ctk.CTkTextbox(self, wrap="word")
        self.html_text.grid(row=0, column=0, padx=16, pady=(16, 8), sticky="nsew")

        # Bottom bar frame
        bottom = ctk.CTkFrame(self)
        bottom.grid(row=1, column=0, padx=16, pady=(0, 16), sticky="ew")
        bottom.grid_columnconfigure(0, weight=1)
        bottom.grid_columnconfigure(1, weight=0)
        bottom.grid_columnconfigure(2, weight=0)
        bottom.grid_columnconfigure(3, weight=0)
        bottom.grid_columnconfigure(4, weight=0)
        bottom.grid_columnconfigure(5, weight=0)
        bottom.grid_columnconfigure(6, weight=0)
        bottom.grid_columnconfigure(7, weight=0)

        # Status label
        self.status_var = ctk.StringVar(value="Ready")
        self.status_label = ctk.CTkLabel(bottom, textvariable=self.status_var, anchor="w")
        self.status_label.grid(row=0, column=0, padx=(12, 12), pady=12, sticky="w")

        # File actions
        self.open_btn = ctk.CTkButton(bottom, text="Open HTML", command=self.on_open_click)
        self.open_btn.grid(row=0, column=1, padx=12, pady=12, sticky="e")

        self.save_btn = ctk.CTkButton(bottom, text="Save HTML", command=self.on_save_click)
        self.save_btn.grid(row=0, column=2, padx=12, pady=12, sticky="e")

        # Paging mode toggle
        self.paging_var = ctk.StringVar(value="Pages")
        self.paging_toggle = ctk.CTkSegmentedButton(
            bottom,
            values=["Pages", "Continuous"],
            variable=self.paging_var,
        )
        self.paging_toggle.grid(row=0, column=3, padx=12, pady=12, sticky="e")

        # Preview button
        self.preview_btn = ctk.CTkButton(bottom, text="Preview HTML", command=self.on_preview_click)
        self.preview_btn.grid(row=0, column=4, padx=12, pady=12, sticky="e")

        # Convert buttons
        self.convert_btn = ctk.CTkButton(bottom, text="Convert to PDF", command=self.on_convert_click)
        self.convert_btn.grid(row=0, column=5, padx=12, pady=12, sticky="e")

        self.convert_docx_btn = ctk.CTkButton(bottom, text="Convert to DOCX", command=self.on_convert_docx_click)
        self.convert_docx_btn.grid(row=0, column=6, padx=12, pady=12, sticky="e")

        self.convert_pptx_btn = ctk.CTkButton(bottom, text="Convert to PPTX", command=self.on_convert_pptx_click)
        self.convert_pptx_btn.grid(row=0, column=7, padx=12, pady=12, sticky="e")

        # Example placeholder
        self._insert_example_placeholder()

        # Live preview state
        self._latest_html: str = self.html_text.get("1.0", "end-1c")
        self._preview_server: Optional[Tuple[socketserver.TCPServer, int, threading.Thread]] = None
        self._debounce_job: Optional[str] = None
        self._highlight_job: Optional[str] = None
        self._current_file: Optional[str] = None

        # Debounced change binding for live preview
        self.html_text.bind("<<Modified>>", self._on_text_modified)
        self.protocol("WM_DELETE_WINDOW", self._on_close)

        # Configure syntax highlight tags
        self._configure_highlight_tags()
        # Load last session if any
        self._load_last_session()
        # Initial highlight
        self._schedule_highlight()

    def _insert_example_placeholder(self) -> None:
        placeholder = (
            "<!doctype html>\n"
            "<html>\n<head>\n<meta charset=\"utf-8\">\n"
            "<meta name=\"viewport\" content=\"width=device-width, initial-scale=1\">\n"
            "<title>Sample</title>\n"
            "<link rel=\"stylesheet\" href=\"https://cdn.jsdelivr.net/npm/tailwindcss@2.2.19/dist/tailwind.min.css\">\n"
            "</head>\n<body class=\"p-12\">\n"
            "<div class=\"prose\">\n"
            "  <h1 class=\"text-3xl font-bold\">HTML-to-PDF Converter</h1>\n"
            "  <p>Paste your HTML here or overwrite this example.</p>\n"
            "  <p><em>External CSS and images will be loaded.</em></p>\n"
            "</div>\n"
            "</body>\n</html>\n"
        )
        self.html_text.insert("1.0", placeholder)

    def _session_path(self) -> str:
        if sys.platform == "win32":
            base = os.path.join(os.environ.get("APPDATA", ""), "HTML-to-PDF Converter")
        elif sys.platform == "darwin":
            base = os.path.expanduser("~/Library/Application Support/HTML-to-PDF Converter")
        else:
            base = os.path.expanduser("~/.config/HTML-to-PDF Converter")
        os.makedirs(base, exist_ok=True)
        return os.path.join(base, "last_session.html")

    def _load_last_session(self) -> None:
        try:
            path = self._session_path()
            if os.path.isfile(path):
                with open(path, "r", encoding="utf-8") as f:
                    data = f.read()
                if data.strip():
                    self.html_text.delete("1.0", "end")
                    self.html_text.insert("1.0", data)
                    self._latest_html = data
                    self.status_var.set("Restored last session")
        except Exception:
            pass

    def _autosave_session(self) -> None:
        try:
            with open(self._session_path(), "w", encoding="utf-8") as f:
                f.write(self._latest_html)
        except Exception:
            pass

    def _configure_highlight_tags(self) -> None:
        # Colors tuned for both light/dark backgrounds
        self.html_text.tag_config("html-comment", foreground="#6b7280")
        self.html_text.tag_config("html-tag", foreground="#2563eb")
        self.html_text.tag_config("html-attr", foreground="#d97706")
        self.html_text.tag_config("html-string", foreground="#16a34a")

    def _schedule_highlight(self) -> None:
        if self._highlight_job is not None:
            try:
                self.after_cancel(self._highlight_job)
            except Exception:
                pass
        self._highlight_job = self.after(150, self._apply_syntax_highlighting)

    def _apply_syntax_highlighting(self) -> None:
        text_widget = self.html_text
        content = text_widget.get("1.0", "end-1c")

        # Clear previous tags
        for tag in ("html-comment", "html-tag", "html-attr", "html-string"):
            text_widget.tag_remove(tag, "1.0", "end")

        # Helper to convert absolute offset to Tk index
        def to_index(offset: int) -> str:
            # Compute line/column by counting newlines
            upto = content[:offset]
            line = upto.count("\n") + 1
            col = len(upto) - (upto.rfind("\n") + 1 if "\n" in upto else 0)
            return f"{line}.{col}"

        # Comments
        for m in re.finditer(r"<!--[\s\S]*?-->", content):
            text_widget.tag_add("html-comment", to_index(m.start()), to_index(m.end()))

        # Strings inside tags
        for m in re.finditer(r"(['\"]).*?\1", content, flags=re.DOTALL):
            text_widget.tag_add("html-string", to_index(m.start()), to_index(m.end()))

        # Tag names and attribute names within <...>
        for m in re.finditer(r"<\s*\/??\s*([a-zA-Z][a-zA-Z0-9:-]*)", content):
            text_widget.tag_add("html-tag", to_index(m.start(1)), to_index(m.end(1)))

        for m in re.finditer(r"\s([a-zA-Z_:][-a-zA-Z0-9_:.]*)\s*=", content):
            text_widget.tag_add("html-attr", to_index(m.start(1)), to_index(m.end(1)))

    def on_convert_click(self) -> None:
        html = self.html_text.get("1.0", "end-1c").strip()
        if not html:
            messagebox.showinfo("No HTML", "Please paste HTML content before converting.")
            return

        output_path = filedialog.asksaveasfilename(
            title="Save PDF As...",
            defaultextension=".pdf",
            filetypes=[("PDF Files", "*.pdf")],
            initialfile="output.pdf",
        )

        if not output_path:
            return

        # Disable button and update status
        self.convert_btn.configure(state="disabled")
        self.status_var.set("Converting...")

        def worker() -> None:
            error_msg: Optional[str] = None
            try:
                continuous = self.paging_var.get() == "Continuous"
                convert_html_to_pdf_sync(html, output_path, continuous=continuous)
            except Exception:
                error_msg = traceback.format_exc()

            def finalize() -> None:
                if error_msg is None:
                    self.status_var.set("PDF saved successfully!")
                else:
                    self.status_var.set("Conversion failed. See details in alert.")
                    messagebox.showerror("Error", f"An error occurred during conversion:\n\n{error_msg}")
                self.convert_btn.configure(state="normal")

            self.after(0, finalize)

        threading.Thread(target=worker, daemon=True).start()

    def on_convert_docx_click(self) -> None:
        html = self.html_text.get("1.0", "end-1c").strip()
        if not html:
            messagebox.showinfo("No HTML", "Please paste HTML content before converting.")
            return

        output_path = filedialog.asksaveasfilename(
            title="Save DOCX As...",
            defaultextension=".docx",
            filetypes=[("Word Document", "*.docx")],
            initialfile="output.docx",
        )

        if not output_path:
            return

        self.convert_docx_btn.configure(state="disabled")
        self.status_var.set("Converting to DOCX...")

        def worker() -> None:
            error_msg: Optional[str] = None
            try:
                convert_html_to_docx_sync(html, output_path)
            except Exception:
                error_msg = traceback.format_exc()

            def finalize() -> None:
                if error_msg is None:
                    self.status_var.set("DOCX saved successfully!")
                else:
                    self.status_var.set("DOCX conversion failed. See details in alert.")
                    messagebox.showerror("Error", f"An error occurred during conversion:\n\n{error_msg}")
                self.convert_docx_btn.configure(state="normal")

            self.after(0, finalize)

        threading.Thread(target=worker, daemon=True).start()

    def on_convert_pptx_click(self) -> None:
        html = self.html_text.get("1.0", "end-1c").strip()
        if not html:
            messagebox.showinfo("No HTML", "Please paste HTML content before converting.")
            return

        output_path = filedialog.asksaveasfilename(
            title="Save PPTX As...",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            initialfile="output.pptx",
        )

        if not output_path:
            return

        self.convert_pptx_btn.configure(state="disabled")
        self.status_var.set("Converting to PPTX...")

        def worker() -> None:
            error_msg: Optional[str] = None
            try:
                convert_html_to_pptx_sync(html, output_path)
            except Exception:
                error_msg = traceback.format_exc()

            def finalize() -> None:
                if error_msg is None:
                    self.status_var.set("PPTX saved successfully!")
                else:
                    self.status_var.set("PPTX conversion failed. See details in alert.")
                    messagebox.showerror("Error", f"An error occurred during conversion:\n\n{error_msg}")
                self.convert_pptx_btn.configure(state="normal")

            self.after(0, finalize)

        threading.Thread(target=worker, daemon=True).start()

    def on_preview_click(self) -> None:
        html = self.html_text.get("1.0", "end-1c").strip()
        if not html:
            messagebox.showinfo("No HTML", "Please paste HTML content to preview.")
            return
        # Start live preview server (once) and open browser
        try:
            if self._preview_server is None:
                self._preview_server = self._start_preview_server()
            _, port, _ = self._preview_server
            webbrowser.open(f"http://127.0.0.1:{port}/")
            self.status_var.set("Opened live preview in browser (auto-refresh)")
        except Exception as exc:
            self.status_var.set("Failed to open preview")
            messagebox.showerror("Error", f"Could not open preview:\n\n{exc}")

    def _on_text_modified(self, _event=None) -> None:
        # Reset modified flag immediately
        try:
            self.html_text.edit_modified(False)
        except Exception:
            pass

        # Debounce updates to reduce churn
        if self._debounce_job is not None:
            try:
                self.after_cancel(self._debounce_job)
            except Exception:
                pass
        self._debounce_job = self.after(300, self._update_latest_html_from_editor)
        self._schedule_highlight()

    def _update_latest_html_from_editor(self) -> None:
        self._latest_html = self.html_text.get("1.0", "end-1c")
        # No UI update necessary; browser polls the server
        self._autosave_session()

    def _start_preview_server(self) -> Tuple[socketserver.TCPServer, int, threading.Thread]:
        app_ref = self

        class Handler(http.server.BaseHTTPRequestHandler):
            def log_message(self, format: str, *args) -> None:  # silence
                return

            def do_GET(self):  # type: ignore[override]
                path = self.path.split("?")[0]
                if path == "/" or path == "/index.html":
                    content = (
                        "<!doctype html>\n"
                        "<html><head><meta charset=\"utf-8\">"
                        "<meta name=\"viewport\" content=\"width=device-width, initial-scale=1\">"
                        "<title>Live Preview</title>"
                        "<style>html,body{height:100%;margin:0}body{background:#0b0b0b}"
                        "#wrap{position:fixed;inset:12px;background:#fff;border-radius:8px;overflow:hidden}"
                        "#bar{position:absolute;top:0;left:0;right:0;height:36px;display:flex;align-items:center;gap:8px;padding:0 12px;background:#f5f5f5;border-bottom:1px solid #e5e5e5;font-family:ui-sans-serif,system-ui,sans-serif}"
                        "#frame{position:absolute;top:36px;left:0;right:0;bottom:0;border:0;width:100%;height:calc(100% - 36px)}"
                        "button{appearance:none;border:1px solid #d0d0d0;border-radius:6px;background:#fff;padding:6px 10px;cursor:pointer}"
                        "</style></head><body>"
                        "<div id=wrap>"
                        "  <div id=bar>Live Preview <button id=refresh>Refresh</button><span id=info style=\"margin-left:auto;color:#666\"></span></div>"
                        "  <iframe id=frame src=/content></iframe>"
                        "</div>"
                        "<script>"
                        "const frame=document.getElementById('frame');"
                        "const info=document.getElementById('info');"
                        "document.getElementById('refresh').onclick=()=>reload();"
                        "function reload(){frame.src='/content?t='+Date.now();info.textContent=new Date().toLocaleTimeString();}"
                        "</script>"
                        "</body></html>\n"
                    ).encode("utf-8")
                    self.send_response(200)
                    self.send_header("Content-Type", "text/html; charset=utf-8")
                    self.send_header("Content-Length", str(len(content)))
                    self.end_headers()
                    self.wfile.write(content)
                elif path == "/content":
                    data = app_ref._latest_html.encode("utf-8")
                    self.send_response(200)
                    self.send_header("Content-Type", "text/html; charset=utf-8")
                    self.send_header("Cache-Control", "no-store")
                    self.send_header("Content-Length", str(len(data)))
                    self.end_headers()
                    self.wfile.write(data)
                else:
                    self.send_error(404)

        # Bind to a free port on localhost
        httpd = socketserver.TCPServer(("127.0.0.1", 0), Handler)
        port = httpd.server_address[1]
        server_thread = threading.Thread(target=httpd.serve_forever, daemon=True)
        server_thread.start()
        return httpd, port, server_thread

    def _on_close(self) -> None:
        # Stop preview server if running
        if self._preview_server is not None:
            try:
                httpd, _, _ = self._preview_server
                httpd.shutdown()
                httpd.server_close()
            except Exception:
                pass
            self._preview_server = None
        # Final autosave
        try:
            self._latest_html = self.html_text.get("1.0", "end-1c")
            self._autosave_session()
        except Exception:
            pass
        self.destroy()

    # ---------- File operations ----------
    def on_open_click(self) -> None:
        path = filedialog.askopenfilename(
            title="Open HTML File",
            filetypes=[("HTML Files", "*.html *.htm"), ("All Files", "*.*")],
        )
        if not path:
            return
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = f.read()
            self.html_text.delete("1.0", "end")
            self.html_text.insert("1.0", data)
            self._latest_html = data
            self._current_file = path
            self.status_var.set(f"Opened: {os.path.basename(path)}")
            self._schedule_highlight()
        except Exception as exc:
            messagebox.showerror("Error", f"Could not open file:\n\n{exc}")

    def on_save_click(self) -> None:
        if self._current_file is None:
            return self.on_save_as_click()
        try:
            data = self.html_text.get("1.0", "end-1c")
            with open(self._current_file, "w", encoding="utf-8") as f:
                f.write(data)
            self.status_var.set(f"Saved: {os.path.basename(self._current_file)}")
        except Exception as exc:
            messagebox.showerror("Error", f"Could not save file:\n\n{exc}")

    def on_save_as_click(self) -> None:
        path = filedialog.asksaveasfilename(
            title="Save HTML As...",
            defaultextension=".html",
            filetypes=[("HTML Files", "*.html"), ("All Files", "*.*")],
            initialfile="document.html",
        )
        if not path:
            return
        try:
            data = self.html_text.get("1.0", "end-1c")
            with open(path, "w", encoding="utf-8") as f:
                f.write(data)
            self._current_file = path
            self.status_var.set(f"Saved: {os.path.basename(path)}")
        except Exception as exc:
            messagebox.showerror("Error", f"Could not save file:\n\n{exc}")


def main() -> None:
    # Only create GUI if we have a display (not in headless environments)
    import os
    import sys

    # Check for display environment
    has_display = False
    if sys.platform == "win32":
        has_display = True  # Windows always has display
    elif sys.platform == "darwin":
        has_display = True  # macOS always has display
    else:
        # Linux/Unix: check for DISPLAY or WAYLAND_DISPLAY
        has_display = bool(os.environ.get("DISPLAY") or os.environ.get("WAYLAND_DISPLAY"))

    if has_display:
        try:
            app = HtmlToPdfApp()
            app.mainloop()
        except Exception as e:
            print(f"GUI Error: {e}")
            print("This might be a headless environment. Use command line conversion instead.")
            sys.exit(1)
    else:
        print("No display detected. This is a GUI application that requires a display.")
        print("Usage: python html_to_pdf_app.py")
        print("Or run in an environment with a display server.")
        sys.exit(1)


if __name__ == "__main__":
    # Handle command line arguments for PyInstaller
    import sys
    if len(sys.argv) > 1 and sys.argv[1] == "--help":
        print("HTML-to-PDF Converter")
        print("A cross-platform GUI app for converting HTML to PDF, DOCX, and PPTX.")
        print("\nUsage:")
        print("  python html_to_pdf_app.py    # Start GUI")
        print("  python html_to_pdf_app.py --help  # Show this help")
        sys.exit(0)

    try:
        main()
    except KeyboardInterrupt:
        print("\nApplication interrupted by user")
        sys.exit(0)
    except Exception as e:
        print(f"Application error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


